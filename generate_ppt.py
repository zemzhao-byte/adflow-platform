"""
自动化投放平台 - 产品架构 PPT 生成脚本 v2
设计原则：大留白、强层次、重点突出、一页一主题
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
W, H = Inches(13.333), Inches(7.5)
prs.slide_width = W
prs.slide_height = H

# ── Design Tokens ──
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BG     = RGBColor(0xF5, 0xF7, 0xFA)
DARK   = RGBColor(0x1B, 0x1F, 0x2B)
BODY   = RGBColor(0x3D, 0x42, 0x50)
MUTED  = RGBColor(0x8C, 0x91, 0x9A)
BORDER = RGBColor(0xE2, 0xE5, 0xEB)
ACCENT = RGBColor(0x3B, 0x82, 0xF6)
GREEN  = RGBColor(0x16, 0xA3, 0x4A)
CYAN   = RGBColor(0x06, 0xB6, 0xD4)
ORANGE = RGBColor(0xEA, 0x58, 0x0C)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
PINK   = RGBColor(0xDB, 0x27, 0x77)
SLATE  = RGBColor(0x64, 0x74, 0x8B)
RED    = RGBColor(0xEF, 0x44, 0x44)
META_C = RGBColor(0x18, 0x77, 0xF2)
GOOG_C = RGBColor(0x42, 0x85, 0xF4)
TT_C   = RGBColor(0x25, 0xF4, 0xEE)
AL_C   = RGBColor(0x19, 0x4E, 0x8B)
FONT   = 'Microsoft YaHei'

def hex_rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def set_bg(slide, color=BG):
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = color

def _tf(shape):
    tf = shape.text_frame; tf.word_wrap = True; return tf

def add_rect(slide, l, t, w, h, fill, radius=True):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                                Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background(); s.shadow.inherit = False
    return s

def label(slide, l, t, w, h, txt, sz, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = _tf(box); p = tf.paragraphs[0]
    p.text = txt; p.font.size = Pt(sz); p.font.color.rgb = color
    p.font.bold = bold; p.font.name = FONT; p.alignment = align
    return box

def tag(slide, l, t, txt, bg_color, fg=WHITE, sz=11, pw=None):
    tw = pw or (len(txt) * 0.11 + 0.35)
    s = add_rect(slide, l, t, tw, 0.32, bg_color)
    tf = _tf(s); p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    p.text = txt; p.font.size = Pt(sz); p.font.color.rgb = fg
    p.font.bold = True; p.font.name = FONT
    return s

def bullets(slide, l, t, w, items, sz=11, color=BODY, spacing=5, leading='•'):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(5))
    tf = _tf(box)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f'{leading} {item}'; p.font.size = Pt(sz)
        p.font.color.rgb = color; p.font.name = FONT
        p.space_before = Pt(spacing); p.space_after = Pt(0)
        p.line_spacing = Pt(sz + 6)
    return box

def section_header(slide, num, title, subtitle=''):
    set_bg(slide, DARK)
    add_rect(slide, 0, 0, 13.333, 7.5, DARK, False)
    tag(slide, 1.2, 2.5, f'PART {num}', ACCENT, WHITE, 14, 1.5)
    label(slide, 1.2, 3.1, 10, 0.8, title, 36, WHITE, True)
    if subtitle:
        label(slide, 1.2, 4.0, 10, 0.5, subtitle, 15, MUTED)

def page_title(slide, title, subtitle=''):
    set_bg(slide)
    add_rect(slide, 0, 0, 13.333, 0.06, ACCENT, False)
    label(slide, 0.9, 0.4, 8, 0.55, title, 24, DARK, True)
    if subtitle:
        label(slide, 0.9, 0.9, 10, 0.35, subtitle, 12, MUTED)


# ═══════════════════════════════════════
# 1. COVER
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK)
add_rect(s, 0, 0, 13.333, 7.5, DARK, False)
add_rect(s, 0, 0, 0.12, 7.5, ACCENT, False)
label(s, 1.2, 1.8, 10, 0.6, 'AdFlow', 18, ACCENT, True)
label(s, 1.2, 2.4, 10, 1.0, '自动化投放平台\n产品架构', 44, WHITE, True)
add_rect(s, 1.2, 3.9, 4, 0.01, ACCENT, False)
label(s, 1.2, 4.2, 10, 0.4, 'Meta · Google Ads · TikTok · AppLovin', 16, RGBColor(0x94,0xA3,0xB8))
label(s, 1.2, 4.7, 10, 0.4, '从账户接入到智能优化的全链路解决方案', 14, MUTED)
label(s, 1.2, 6.3, 4, 0.3, '2026  ·  产品团队', 12, SLATE)


# ═══════════════════════════════════════
# 2. TOC
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
page_title(s, '目录')

toc = [
    ('01', '海外投放全链路', '6大阶段完整闭环'),
    ('02', '六层系统架构', '分层设计 · 模块拆解 · 功能清单'),
    ('03', '建设节奏', '6个Phase · 分阶段交付'),
    ('04', '平台能力对比', '四大渠道 7 维度矩阵'),
    ('05', 'Meta 配置链路', 'Campaign→Ad Set→Ad 全字段'),
    ('06', 'Google 配置链路', 'UAC · PMax · Smart Bidding'),
    ('07', '归因与数据', 'AppsFlyer · Adjust · SKAN · 融合'),
    ('08', '业务能力', '素材 · 广告 · 报表 · 规则'),
    ('09', '基础设施', '权限 · 模板 · 监控 · 安全'),
]
colors = [ACCENT, GREEN, CYAN, ORANGE, META_C, GOOG_C, ORANGE, CYAN, SLATE]
for i, (num, title, desc) in enumerate(toc):
    row = i // 3; col = i % 3
    x = 0.9 + col * 4.0; y = 1.6 + row * 1.9
    card = add_rect(s, x, y, 3.6, 1.5, WHITE)
    card.line.color.rgb = BORDER; card.line.width = Pt(0.75)
    add_rect(s, x, y, 0.08, 1.5, colors[i], False)
    label(s, x + 0.3, y + 0.2, 1, 0.4, num, 28, colors[i], True)
    label(s, x + 0.3, y + 0.65, 3, 0.35, title, 15, DARK, True)
    label(s, x + 0.3, y + 1.0, 3, 0.3, desc, 10, MUTED)


# ═══════════════════════════════════════
# 3. SECTION: 投放全链路
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
section_header(s, '01', '海外投放全链路', '从投放准备到数据复盘的 6 大阶段闭环')

s = prs.slides.add_slide(prs.slide_layouts[6])
page_title(s, '海外投放全链路', '每个阶段对应平台的核心能力模块')

stages = [
    ('投放准备', ACCENT, '开户·SDK·埋点·素材', ['平台开户(BM/BC)', 'MMP SDK接入', '事件埋点配置', '追踪链接生成', '素材规格制作']),
    ('广告创编', CYAN, 'Campaign三层搭建', ['选择投放目标', 'Campaign配置', 'Ad Set定向出价', 'Ad素材绑定', 'Dry-run校验']),
    ('投放上线', RGBColor(0xCA,0x8A,0x04), '审核→学习→放量', ['平台审核', 'Learning Phase', '预算爬坡', '多平台同步上线']),
    ('投放优化', GREEN, '实时数据驱动调优', ['实时监控', '预算调控', '出价调优', '素材轮替', '定向拓展']),
    ('自动化规则', PURPLE, '7×24自动止损加量', ['止损规则', '加量规则', '素材替换', '异常告警', '定时调度']),
    ('数据分析', ORANGE, '归因·LTV·复盘', ['实时报表', '归因分析', 'Cohort/LTV', '跨渠道对比', '策略复盘']),
]
for i, (name, color, sub, items) in enumerate(stages):
    x = 0.4 + i * 2.1
    card = add_rect(s, x, 1.4, 1.95, 5.5, WHITE)
    card.line.color.rgb = BORDER; card.line.width = Pt(0.5)
    add_rect(s, x, 1.4, 1.95, 0.06, color, False)
    label(s, x + 0.15, 1.65, 1.7, 0.4, name, 16, DARK, True)
    label(s, x + 0.15, 2.05, 1.7, 0.25, sub, 9, MUTED)
    add_rect(s, x + 0.15, 2.4, 1.65, 0.01, BORDER, False)
    bullets(s, x + 0.15, 2.55, 1.7, items, 10, BODY, 3)
    if i < 5:
        label(s, x + 1.95, 3.5, 0.2, 0.4, '→', 16, MUTED, False, PP_ALIGN.CENTER)


# ═══════════════════════════════════════
# 4. SECTION: 系统架构
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
section_header(s, '02', '六层系统架构', '从账户接入到数据复盘的完整分层设计')

arch_layers = [
    ('账户与资产管理', GREEN, '投放的基础层', [
        ('平台账户管理', '多渠道账户绑定 · Token安全存储 · 增删改查'),
        ('应用管理', 'App注册 · MMP SDK检测 · 事件埋点清单'),
        ('渠道返点管理', '代理商维护 · 固定/阶梯返点 · 结算对账'),
    ]),
    ('广告投放层', CYAN, '优化师核心操作', [
        ('广告创编', '6步向导 · 渠道自适配 · Dry-run校验'),
        ('素材管理', '上传校验 · 标签分类 · 效果看板 · 疲劳监控'),
        ('广告模板', '一键复用 · 定向模板 · 素材组模板'),
        ('广告运营', '批量启停 · 预算调整 · 爬坡策略'),
    ]),
    ('平台对接层', ACCENT, '屏蔽平台差异', [
        ('统一Connector', 'BaseConnector抽象 · 4渠道实现 · 错误码映射'),
        ('字段映射引擎', '统一模型↔平台API双向映射'),
        ('API弹性机制', '限流 · 重试 · 幂等 · 熔断'),
    ]),
    ('数据与归因层', ORANGE, 'MMP对接与报表', [
        ('归因数据接入', 'AppsFlyer/Adjust · Push/Pull · SKAN 4.0'),
        ('数据报表', '实时KPI · 多维报表 · Cohort/LTV'),
        ('数据仓库', 'PostgreSQL · ClickHouse · ETL'),
    ]),
    ('智能化层', PURPLE, '替代人工盯盘', [
        ('规则引擎', '止损/加量/疲劳 · 冷却期 · 审计日志'),
        ('智能预算分配', 'CPI/ROAS驱动 · 跨渠道优化'),
        ('异常检测', '花费突变 · CPI暴涨 · Learning保护'),
    ]),
    ('基础设施层', SLATE, '安全与运维', [
        ('用户与权限', '5角色 · 审计 · 多租户'),
        ('任务调度', '定时 · 异步队列 · 重试'),
        ('监控告警', 'API监控 · 限流告警 · 健康检查'),
    ]),
]
for page in range(2):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    page_title(s, f'系统架构全景（{"上" if page==0 else "下"}）')
    group = arch_layers[page*3:(page+1)*3]
    for li, (name, color, desc, modules) in enumerate(group):
        y = 1.4 + li * 2.0
        tag(s, 0.6, y, name, color, WHITE, 13, 2.0)
        label(s, 2.8, y + 0.02, 3, 0.3, desc, 11, MUTED)
        for mi, (mname, mdesc) in enumerate(modules):
            x = 0.6 + mi * 3.15
            card = add_rect(s, x, y + 0.5, 3.0, 1.1, WHITE)
            card.line.color.rgb = BORDER; card.line.width = Pt(0.5)
            add_rect(s, x, y + 0.5, 3.0, 0.04, color, False)
            label(s, x + 0.2, y + 0.65, 2.6, 0.3, mname, 13, DARK, True)
            label(s, x + 0.2, y + 0.95, 2.6, 0.5, mdesc, 9, MUTED)
        # data flow arrow
        if li < len(group) - 1:
            label(s, 6.3, y + 1.65, 0.5, 0.3, '↓', 16, MUTED, False, PP_ALIGN.CENTER)


# ═══════════════════════════════════════
# 5. 建设节奏
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
section_header(s, '03', '建设节奏与优先级', '每个阶段有明确交付目标，确保每步可验证')

s = prs.slides.add_slide(prs.slide_layouts[6])
page_title(s, '分阶段建设路线图')
phases = [
    ('Phase 1', '地基搭建', '2-3周', GREEN, '平台账户管理 · 应用管理 · 统一Connector(Meta+Google)', '能绑定Meta和Google账户，Token安全存储，基础API调通'),
    ('Phase 2', '核心投放', '3-4周', CYAN, '广告创编 · 广告模板 · 字段映射引擎 · 素材管理', '能通过平台完成Meta+Google的完整创编流程'),
    ('Phase 3', '数据闭环', '2-3周', ORANGE, '归因数据接入(AF) · 数据报表 · 数据仓库', '完整投放数据报表，花费-安装-CPI闭环'),
    ('Phase 4', '运营效率', '2-3周', PURPLE, '广告运营(批量) · 规则引擎 · 异常检测 · 渠道返点', '规则引擎替代人工盯盘，运营效率工具'),
    ('Phase 5', '智能化', '3-4周', PINK, '智能预算分配 · Cohort/LTV · A/B Test · TT+AL Connector', '智能化自动投放，数据驱动决策闭环'),
    ('Phase 6', '企业级', '2-3周', SLATE, '用户与权限 · 多租户 · 监控告警 · 任务调度', '多团队协作，安全合规，生产级稳定性'),
]
for i, (phase, name, dur, color, modules, goal) in enumerate(phases):
    y = 1.3 + i * 1.0
    # timeline dot + line
    cx = 1.0
    add_rect(s, cx - 0.05, y + 0.08, 0.25, 0.25, color)
    if i < 5:
        add_rect(s, cx + 0.06, y + 0.35, 0.03, 0.65, BORDER, False)
    label(s, 1.4, y, 2.5, 0.35, f'{phase}: {name}', 15, DARK, True)
    tag(s, 3.9, y + 0.03, dur, RGBColor(0xF1,0xF5,0xF9), SLATE, 10, 0.8)
    label(s, 5.0, y, 3.8, 0.35, modules, 9, MUTED)
    card = add_rect(s, 9.0, y - 0.05, 3.8, 0.45, RGBColor(0xF0,0xFD,0xF4))
    card.line.fill.background()
    label(s, 9.2, y, 3.5, 0.35, '🎯 ' + goal, 10, GREEN)


# ═══════════════════════════════════════
# 6. 平台能力对比
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
section_header(s, '04', '广告平台能力对比', 'Meta · Google Ads · TikTok · AppLovin  七维度矩阵')

s = prs.slides.add_slide(prs.slide_layouts[6])
page_title(s, '四大平台能力对比矩阵')

dims = ['API版本', '广告结构', '最低日预算', '出价策略', '素材类型', '数据延迟', 'API限流']
plats = [
    ('Meta', META_C, ['v21.0', 'Campaign→Ad Set→Ad', '$1', 'Lowest Cost / Cost Cap\nBid Cap / Min ROAS', '图片/视频/轮播\nPlayable', '15-30min', '200写/小时/账户']),
    ('Google', GOOG_C, ['v18', 'Campaign→Ad Group\n→Ad/Asset Group', '$1 (建议≥$50)', 'Max Conv / Target CPA\nTarget ROAS', '图片/视频/HTML5\nApp开屏', '3-6h', '15,000请求/天']),
    ('TikTok', RGBColor(0x00,0x00,0x00), ['v1.3', 'Campaign→Ad Group→Ad', '$20', 'Lowest Cost\nCost Cap', '视频/图片/Playable\nSpark Ads', '1-3h', '10 QPS/App']),
    ('AppLovin', AL_C, ['Mgmt API v1', 'Campaign→Ad Group→Ad', '$100', 'Target CPI/CPA\nTarget ROAS', '视频/Playable/Banner\nRewarded', '2-4h', '~5 QPS建议']),
]

def slide_val_label(slide, l, t, w, txt):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(0.5))
    tf = _tf(box); p = tf.paragraphs[0]
    p.text = txt; p.font.size = Pt(9); p.font.color.rgb = BODY
    p.font.name = FONT; p.alignment = PP_ALIGN.CENTER
    p.line_spacing = Pt(13)
    return box

col_start = 2.6; col_w = 2.6; gap = 0.1
for pi, (pn, pc, _) in enumerate(plats):
    x = col_start + pi * (col_w + gap)
    tag(s, x, 1.3, pn, pc, WHITE, 13, col_w)

for di, dim in enumerate(dims):
    y = 1.85 + di * 0.73
    bg = WHITE if di % 2 == 0 else BG
    r = add_rect(s, 0.5, y, 2.0, 0.65, bg); r.line.fill.background()
    label(s, 0.7, y + 0.15, 1.8, 0.4, dim, 11, DARK, True)
    for pi, (_, _, vals) in enumerate(plats):
        x = col_start + pi * (col_w + gap)
        r = add_rect(s, x, y, col_w, 0.65, bg); r.line.fill.background()
        slide_val_label(s, x + 0.15, y + 0.08, col_w - 0.3, vals[di])


# ═══════════════════════════════════════
# 7. Meta 配置链路
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
section_header(s, '05', 'Meta 投放配置链路', 'Campaign → Ad Set → Ad  完整字段与 Advantage+ 智能功能')

# Campaign层
s = prs.slides.add_slide(prs.slide_layouts[6])
page_title(s, 'Meta · Campaign 层配置', 'v21.0+  ·  所有配置的起点')
fields = [
    ('objective *', '投放目标', 'OUTCOME_APP_PROMOTION / TRAFFIC / ENGAGEMENT / LEADS / SALES / AWARENESS', '必填，v21.0起统一OUTCOME_前缀'),
    ('special_ad_categories *', '特殊广告类别', 'NONE / HOUSING / EMPLOYMENT / CREDIT / POLITICS', '必填，非NONE时定向受限'),
    ('CBO', 'Campaign预算优化', '开关', '开启后预算在Campaign层统一管理，系统自动分配给各Ad Set'),
    ('bid_strategy', '出价策略', 'LOWEST_COST / COST_CAP / BID_CAP / MINIMUM_ROAS', 'CBO开启时在此设置'),
    ('daily_budget', 'Campaign预算', '金额（最低$1/天）', 'CBO开启时必填，总预算需设结束日期'),
    ('Advantage+ Budget', '智能预算', '自动', 'v25.0起配合Advantage+受众和版位自动进入'),
]
for i, (field, name, options, note) in enumerate(fields):
    y = 1.3 + i * 0.95
    bg = WHITE if i % 2 == 0 else BG
    r = add_rect(s, 0.5, y, 12.3, 0.85, bg); r.line.fill.background()
    label(s, 0.7, y + 0.08, 2.2, 0.3, field, 12, META_C, True)
    label(s, 0.7, y + 0.38, 2.2, 0.3, name, 10, MUTED)
    label(s, 3.1, y + 0.08, 4.5, 0.65, options, 10, DARK)
    label(s, 7.8, y + 0.08, 4.8, 0.65, note, 9, MUTED)

# Ad Set层
s = prs.slides.add_slide(prs.slide_layouts[6])
page_title(s, 'Meta · Ad Set 层配置', '定向 · 版位 · 出价 · 排期 · Advantage+ 受众')
fields2 = [
    ('optimization_goal *', '优化目标', 'APP_INSTALLS / CONVERSIONS / VALUE / LINK_CLICKS / REACH / THRUPLAY', '需与Campaign objective匹配'),
    ('geo_locations *', '地区定向', '国家/地区/城市/ZIP/GPS半径', '必填，Advantage+不扩展此项'),
    ('Advantage+ 受众 ⚡', '智能受众', 'v23.0起默认开启', '定向变为"建议信号"，系统可扩展到建议外用户'),
    ('Advantage+ 版位 ⚡', '智能版位', '全选所有版位（推荐）', 'Meta自动分配预算到最佳版位'),
    ('promoted_object *', '推广对象', 'application_id + object_store_url + event_type', 'App类目标必填'),
    ('attribution_spec', '归因窗口', '默认7天点击+1天曝光', '可配: click_window=1/7/28天, view_window=0/1天'),
    ('bid_amount', '出价金额', '金额', 'COST_CAP/BID_CAP时必填'),
]
for i, (field, name, options, note) in enumerate(fields2):
    y = 1.3 + i * 0.82
    bg = WHITE if i % 2 == 0 else BG
    r = add_rect(s, 0.5, y, 12.3, 0.72, bg); r.line.fill.background()
    is_smart = '⚡' in field
    label(s, 0.7, y + 0.08, 2.5, 0.3, field.replace(' ⚡',''), 12, PURPLE if is_smart else META_C, True)
    if is_smart:
        tag(s, 3.0, y + 0.07, '智能', PURPLE, WHITE, 8, 0.5)
    label(s, 0.7, y + 0.38, 2.5, 0.3, name, 10, MUTED)
    label(s, 3.6, y + 0.08, 4.2, 0.55, options, 10, DARK)
    label(s, 7.9, y + 0.08, 4.8, 0.55, note, 9, MUTED)

# Ad层
s = prs.slides.add_slide(prs.slide_layouts[6])
page_title(s, 'Meta · Ad 层配置 & 智能功能', '素材 · 文案 · CTA · Advantage+ 创意增强')
ad_fields = [
    ('creative *', '广告创意', '需先POST /act_{id}/adcreatives创建Creative对象', '获取creative_id后绑定'),
    ('素材 *', '图片/视频', '图片: adimages→hash · 视频: advideos→video_id', '支持轮播(2-10张)、Playable(≤2MB)'),
    ('文案', '标题/正文', 'title≤40字符 · body≤125字符(建议) · description≤30字符', '自动截断风险需前端校验'),
    ('CTA', '行动按钮', 'INSTALL_MOBILE_APP / LEARN_MORE / SHOP_NOW / PLAY_GAME等', 'App类常用INSTALL_MOBILE_APP'),
    ('url_tags', 'UTM追踪', 'utm_source=facebook&utm_campaign={{campaign.name}}', '自动追加到落地页链接'),
    ('Advantage+ 创意 ⚡', '创意增强', 'v22.0+ 可独立控制各增强项', 'image_templates / text_optimizations / music等'),
]
for i, (field, name, options, note) in enumerate(ad_fields):
    y = 1.3 + i * 0.85
    bg = WHITE if i % 2 == 0 else BG
    r = add_rect(s, 0.5, y, 12.3, 0.75, bg); r.line.fill.background()
    is_smart = '⚡' in field
    label(s, 0.7, y + 0.08, 2.5, 0.3, field.replace(' ⚡',''), 12, PURPLE if is_smart else META_C, True)
    label(s, 0.7, y + 0.38, 2.5, 0.3, name, 10, MUTED)
    label(s, 3.3, y + 0.08, 4.5, 0.55, options, 10, DARK)
    label(s, 7.9, y + 0.08, 4.8, 0.55, note, 9, MUTED)

# Key notes
add_rect(s, 0.5, 6.5, 12.3, 0.04, ORANGE, False)
notes_text = '⚠ v23.0+ advantage_audience默认=1  ·  v25.0 废弃smart_promotion_type  ·  Learning Phase需~50次转化(3-5天)'
label(s, 0.7, 6.6, 12, 0.4, notes_text, 10, ORANGE)


# ═══════════════════════════════════════
# 8. Google 配置链路
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
section_header(s, '06', 'Google Ads 配置链路', 'UAC 全自动化  ·  Performance Max  ·  Smart Bidding')

s = prs.slides.add_slide(prs.slide_layouts[6])
page_title(s, 'Google Ads · Campaign 层', 'API v18  ·  App Campaign 高度自动化')
g_fields = [
    ('channel_type *', '广告系列类型', 'MULTI_CHANNEL(UAC) / PERFORMANCE_MAX / SEARCH / DISPLAY / VIDEO', 'App推广最常用MULTI_CHANNEL'),
    ('channel_sub_type', '子类型', 'APP_CAMPAIGN(安装) / FOR_ENGAGEMENT(互动) / FOR_PRE_REGISTRATION', 'APP_CAMPAIGN最常用'),
    ('app_id *', '应用ID', 'Android: com.example.app · iOS: 123456789', '与app_store对应'),
    ('bidding_goal *', '出价目标', '优化安装CPI / 优化应用内转化 / 目标CPA / 目标ROAS', '【核心】决定Google ML优化方向'),
    ('campaign_budget *', '日预算', '需单独创建CampaignBudget资源', '建议≥$50/天 (≥10×CPI)'),
    ('geo/language', '地区语言', '通过CampaignCriterion添加GeoTarget', 'UAC不支持手动定向'),
    ('network_settings', '网络设置', 'UAC: Google自动分配Search/Display/YouTube/Play', '不可手动选择'),
]
for i, (field, name, options, note) in enumerate(g_fields):
    y = 1.3 + i * 0.82
    bg = WHITE if i % 2 == 0 else BG
    r = add_rect(s, 0.5, y, 12.3, 0.72, bg); r.line.fill.background()
    label(s, 0.7, y + 0.08, 2.2, 0.3, field, 12, GOOG_C, True)
    label(s, 0.7, y + 0.38, 2.2, 0.3, name, 10, MUTED)
    label(s, 3.1, y + 0.08, 4.5, 0.55, options, 10, DARK)
    label(s, 7.8, y + 0.08, 4.8, 0.55, note, 9, MUTED)

s = prs.slides.add_slide(prs.slide_layouts[6])
page_title(s, 'Google Ads · Ad Group & Asset 层', '素材是唯一可控杠杆')

left_items = [
    ('Ad Group 层（UAC下几乎无配置）', GOOG_C, [
        'App Campaign: 全由Google ML控制',
        '不支持手动添加关键词/受众/定向',
        'ENGAGEMENT类型可覆盖Ad Group级CPA',
        'ENGAGEMENT必须添加user_list受众',
    ]),
]
right_items = [
    ('Ad / Asset 层 *', GOOG_C, [
        '标题: 最多5个，每个≤30字符 [必填]',
        '描述: 最多5个，每个≤90字符 [必填]',
        '图片: 推荐≥3张 (1200×628横版必选)',
        '视频: YouTube视频ID，推荐3+个(横/竖/方)',
        'HTML5: Playable广告 ≤1MB [可选]',
        'Google自动组合测试所有素材',
    ]),
]
for items_list, start_x in [(left_items, 0.5), (right_items, 6.5)]:
    for _, color, items in items_list:
        card = add_rect(s, start_x, 1.3, 6.0, 4.5, WHITE)
        card.line.color.rgb = BORDER; card.line.width = Pt(0.5)
        add_rect(s, start_x, 1.3, 6.0, 0.04, color, False)
        label(s, start_x + 0.3, 1.5, 5.4, 0.35, items_list[0][0], 15, DARK, True)
        bullets(s, start_x + 0.3, 2.0, 5.4, items, 11, BODY, 6)

# Key callouts
add_rect(s, 0.5, 6.0, 12.3, 0.04, RED, False)
label(s, 0.7, 6.1, 12, 0.4, '⚠ UAC高度自动化: 版位/定向/出价全ML控制，素材是唯一杠杆  ·  Conversion Tracking必须提前配好  ·  预算≥CPI×10/天', 10, RED)


# ═══════════════════════════════════════
# 9. 归因与数据
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
section_header(s, '07', '归因与数据体系', 'AppsFlyer · Adjust · SKAN · 数据融合')

s = prs.slides.add_slide(prs.slide_layouts[6])
page_title(s, '归因平台对比与数据融合')

# AF card
card = add_rect(s, 0.5, 1.3, 5.9, 3.5, WHITE)
card.line.color.rgb = BORDER; card.line.width = Pt(0.5)
add_rect(s, 0.5, 1.3, 5.9, 0.04, GREEN, False)
label(s, 0.8, 1.5, 4, 0.35, 'AppsFlyer', 18, GREEN, True)
af = ['Click-through归因: 7天窗口', 'View-through: 1天', 'SKAN 4.0: Coarse/Fine Conversion Value',
      'Push API: 实时Callback', 'Pull API: 批量拉取', 'S2S Events: 服务端上报',
      'Protect360反作弊', 'OneLink追踪链接', 'Cohort API: D1-D365 LTV']
bullets(s, 0.8, 2.0, 5.4, af, 10, BODY, 3)

# Adjust card
card = add_rect(s, 6.8, 1.3, 5.9, 3.5, WHITE)
card.line.color.rgb = BORDER; card.line.width = Pt(0.5)
add_rect(s, 6.8, 1.3, 5.9, 0.04, ACCENT, False)
label(s, 7.1, 1.5, 4, 0.35, 'Adjust', 18, ACCENT, True)
adj = ['Click归因: 7天 · Impression: 24h', 'Fingerprint: 24h', 'SKAN 4.0支持',
       'Callback: 实时推送 100+宏参数', 'Pull API + CSV Export + Datascape',
       'Fraud Prevention反作弊', 'GDPR / COPPA合规',
       '内置Partner Module(Meta/TikTok/Google)', 'Tracker URL管理']
bullets(s, 7.1, 2.0, 5.4, adj, 10, BODY, 3)

# Data Fusion
card = add_rect(s, 0.5, 5.1, 12.2, 2.0, WHITE)
card.line.color.rgb = BORDER; card.line.width = Pt(0.5)
add_rect(s, 0.5, 5.1, 12.2, 0.04, ORANGE, False)
label(s, 0.8, 5.3, 4, 0.35, '数据融合逻辑', 15, ORANGE, True)
fusion = [
    '平台Spend(15min-6h延迟) + MMP归因(实时-3h) → 以MMP归因为准计算CPI/ROAS',
    '数据回溯: 归因窗口内持续修正(Meta最长28天)，需回刷机制',
    'SKAN特殊处理: 延迟24-72h · Campaign粒度 · Conversion Value解码映射',
]
bullets(s, 0.8, 5.7, 11.5, fusion, 10, BODY, 4)


# ═══════════════════════════════════════
# 10. 业务能力
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
section_header(s, '08', '业务能力模块', '素材管理 · 广告管理 · 标准数据 · 规则引擎')

s = prs.slides.add_slide(prs.slide_layouts[6])
page_title(s, '核心业务能力')

biz = [
    ('素材管理', ORANGE, [
        '上传+格式校验(分辨率/时长/编码)',
        '标签分类(类型/语言/地区/风格)',
        '智能裁剪 16:9 / 9:16 / 1:1',
        '效果看板(CTR/CVR/CPI/IPM)',
        '疲劳监控(CTR下降→推荐替换)',
        '各平台规格自动适配',
    ]),
    ('广告管理', CYAN, [
        '6步创编向导(完整链路)',
        '跨平台字段自动映射',
        '严格校验+Dry-run',
        '投放目标体系(Install/AEO/VO)',
        '模板复用(Campaign/定向/素材)',
        '批量操作(启停/预算/复制)',
    ]),
    ('标准数据', ORANGE, [
        'KPI: Spend/CPI/CTR/CVR/ROAS',
        '多维报表(时间/平台/地区/素材)',
        'MMP归因融合(以归因为准)',
        'Cohort/LTV(D1-D365)',
        '跨渠道对比',
        '导出(CSV/Excel/API)',
    ]),
    ('规则引擎', PURPLE, [
        '止损: CPI>目标×1.3→暂停',
        '加量: CPI<目标×0.8→预算+20%',
        '零转化: 消耗>$30无Install→停',
        '素材疲劳: CTR连降3天→停',
        '冷却期+审批模式',
        '全链路执行日志审计',
    ]),
]
for i, (name, color, items) in enumerate(biz):
    x = 0.5 + i * 3.15
    card = add_rect(s, x, 1.3, 2.95, 5.8, WHITE)
    card.line.color.rgb = BORDER; card.line.width = Pt(0.5)
    add_rect(s, x, 1.3, 2.95, 0.04, color, False)
    label(s, x + 0.2, 1.5, 2.5, 0.35, name, 16, DARK, True)
    add_rect(s, x + 0.2, 1.95, 2.55, 0.01, BORDER, False)
    bullets(s, x + 0.2, 2.1, 2.55, items, 10, BODY, 4)


# ═══════════════════════════════════════
# 11. 基础设施
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
section_header(s, '09', '基础设施与安全', '权限 · 审计 · 模板 · 监控')

s = prs.slides.add_slide(prs.slide_layouts[6])
page_title(s, '基础设施')

infra = [
    ('用户与权限', SLATE, [
        'Admin: 全局配置·用户管理·绑定',
        'Lead: 数据查看·规则审批·团队',
        'Optimizer: 创编/调整(限本人)',
        'Viewer: 仅查看报表',
        'Finance: 花费/预算/账单',
        '操作审计日志(全写操作)',
        '大额调整需主管审批',
        'Token AES-256加密存储',
    ]),
    ('模板与效率', PINK, [
        'Campaign模板一键复用',
        '定向模板(T1-iOS-Core等)',
        '素材组模板(素材+文案+CTA)',
        '命名规则自动生成',
        '批量启停/调预算',
        '跨平台上线(一键推送)',
        'Excel导入批量创建',
    ]),
    ('任务与监控', ACCENT, [
        '定时任务(报表/规则/Token)',
        '异步队列(批量/大规模创编)',
        '任务重试+死信队列',
        'API成功率/延迟监控',
        '平台限流告警',
        '服务健康检查',
        '业务指标异常告警',
    ]),
    ('技术栈', GREEN, [
        'Backend: FastAPI + Python',
        'DB: PostgreSQL(业务)',
        'Analytics: ClickHouse(时序)',
        'Cache: Redis',
        'Queue: Celery + Temporal',
        'Frontend: React + Ant Design',
        'Monitor: OpenTelemetry',
    ]),
]
for i, (name, color, items) in enumerate(infra):
    x = 0.3 + i * 3.2
    card = add_rect(s, x, 1.3, 3.0, 5.8, WHITE)
    card.line.color.rgb = BORDER; card.line.width = Pt(0.5)
    add_rect(s, x, 1.3, 3.0, 0.04, color, False)
    label(s, x + 0.2, 1.5, 2.6, 0.35, name, 15, DARK, True)
    add_rect(s, x + 0.2, 1.9, 2.6, 0.01, BORDER, False)
    bullets(s, x + 0.2, 2.05, 2.6, items, 10, BODY, 3)


# ═══════════════════════════════════════
# FINAL: Thank you
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK)
add_rect(s, 0, 0, 13.333, 7.5, DARK, False)
add_rect(s, 0, 0, 0.12, 7.5, ACCENT, False)
label(s, 1.2, 2.5, 10, 0.8, 'Thank You', 44, WHITE, True)
add_rect(s, 1.2, 3.5, 3, 0.01, ACCENT, False)
label(s, 1.2, 3.8, 10, 0.5, '自动化投放平台 · 产品架构', 18, RGBColor(0x94,0xA3,0xB8))
label(s, 1.2, 4.5, 10, 0.4, '本文档持续迭代，产品架构更新时同步更新', 13, MUTED)
label(s, 1.2, 5.0, 10, 0.4, '生成脚本: generate_ppt.py', 11, SLATE)

# ── Save ──
out = '/Users/zem.zhao/Documents/自动化投放/自动化投放平台_产品架构.pptx'
prs.save(out)
print(f'✅ PPT v2 已生成: {out}')
print(f'   共 {len(prs.slides)} 页')
